"""PowerPoint parser — slide-by-slide text + notes extraction."""

from __future__ import annotations

import asyncio
from pathlib import Path

from core.models import DocumentType

from .base import BaseParser, ParsedDocument


class PowerPointParser(BaseParser):
    """Parse PowerPoint files (.pptx, .ppt) with slide structure."""

    document_type = DocumentType.POWERPOINT
    supported_mimes = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]

    async def parse(self, file_path: Path) -> ParsedDocument:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._parse_sync, file_path)

    def _parse_sync(self, file_path: Path) -> ParsedDocument:
        import pptx

        prs = pptx.Presentation(str(file_path))

        text_parts = []
        slide_data = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            tables_found = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text_frame.text.strip()
                    if shape_text:
                        texts.append(shape_text)

                # Extract tables from slides
                if shape.has_table:
                    table = shape.table
                    table_text = []
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text for cell in row.cells
                        )
                        table_text.append(row_text)
                    tables_found.append({
                        "rows": len(table.rows),
                        "columns": len(table.columns),
                        "text": "\n".join(table_text),
                    })

            # Notes
            notes = ""
            if slide.has_notes_slide:
                notes = (
                    slide.notes_slide.notes_text_frame.text.strip()
                )

            slide_text = "\n".join(texts)
            if notes:
                slide_text += f"\n[Notes]: {notes}"

            text_parts.append(f"## Slide {i}\n{slide_text}")

            slide_data.append({
                "slide_number": i,
                "shape_count": len(slide.shapes),
                "text_length": len(slide_text),
                "has_notes": bool(notes),
                "table_count": len(tables_found),
            })

        text = "\n\n".join(text_parts)
        metadata = {
            "slide_count": len(prs.slides),
            "slide_data": slide_data,
            "has_notes": any(s["has_notes"] for s in slide_data),
        }

        return ParsedDocument(
            text=text,
            metadata=metadata,
        )

    def can_handle(self, file_path: Path, mime_type: str) -> bool:
        if mime_type in self.supported_mimes:
            return True
        return file_path.suffix.lower() in (".pptx", ".ppt")
